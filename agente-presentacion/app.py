"""UI Streamlit para el Agente Presentacion.

Flujo:
  1. El profesor sube el .md generado por el Agente Contenido.
  2. Pulsa "Detectar elementos" -> se ejecuta detector.detectar_elementos().
  3. La app muestra los elementos detectados con checkboxes.
  4. El profesor selecciona los que quiere exportar.
  5. Pulsa "Generar PDF" -> generador_pdf.generar_pdf() -> descarga.
     Pulsa "Generar HTML interactivo" -> generador_html.generar_html()
"""

from __future__ import annotations

import hashlib
import io
import re

import streamlit as st
import streamlit.components.v1 as components

from detector import detectar_elementos
from generador_html import generar_html, _slug as _slug_html
from generador_pdf import generar_pdf

# ---------------------------------------------------------------------------
# Hero HTML — Agente 03, workflow de 4 pasos
# ---------------------------------------------------------------------------

_HERO_PRES_HTML = r"""<!DOCTYPE html><html><head>
<link href="https://fonts.googleapis.com/css2?family=Playfair+Display:wght@400;500&family=DM+Sans:wght@400;500&display=swap" rel="stylesheet">
<style>
*{margin:0;padding:0;box-sizing:border-box;}
:root{
  --text1:#2C2C2A;--text2:#5F5E5A;--text3:#888780;
  --card:rgba(0,0,0,0.03);--border:rgba(0,0,0,0.1);--arrow:rgba(0,0,0,0.2);
}
:root.dark{
  --text1:#FAFAFA;--text2:rgba(255,255,255,0.65);--text3:rgba(255,255,255,0.4);
  --card:rgba(255,255,255,0.06);--border:rgba(255,255,255,0.1);--arrow:rgba(255,255,255,0.22);
}
@media(prefers-color-scheme:dark){:root:not(.light){
  --text1:#FAFAFA;--text2:rgba(255,255,255,0.65);--text3:rgba(255,255,255,0.4);
  --card:rgba(255,255,255,0.06);--border:rgba(255,255,255,0.1);--arrow:rgba(255,255,255,0.22);
}}
body{background:transparent;font-family:'DM Sans',sans-serif;overflow:hidden;padding:0 2px;}
.hero{padding:32px 0 16px 0;}
.eyebrow{font-size:11px;font-weight:500;color:#185FA5;letter-spacing:.14em;
  text-transform:uppercase;margin-bottom:12px;}
.title{font-family:'Playfair Display',serif;font-size:38px;font-weight:500;
  color:var(--text1);line-height:1.15;margin-bottom:14px;}
.title .accent{color:#185FA5;}
.desc{font-size:15px;font-weight:400;color:var(--text2);line-height:1.6;max-width:560px;}
.workflow{display:flex;align-items:center;margin-top:28px;padding:16px 20px;
  background:var(--card);border:.5px solid var(--border);border-radius:10px;}
.step{display:flex;align-items:center;gap:10px;flex:1;}
.num{width:28px;height:28px;border-radius:50%;background:#185FA5;color:#FFF;
  font-size:12px;font-weight:500;display:flex;align-items:center;justify-content:center;
  flex-shrink:0;box-shadow:0 2px 8px rgba(24,95,165,.25);}
.lbl{font-size:9px;font-weight:500;color:var(--text3);text-transform:uppercase;
  letter-spacing:.08em;margin-bottom:2px;}
.sdesc{font-size:12px;font-weight:500;color:var(--text1);}
.arrow{flex-shrink:0;margin:0 4px;color:var(--arrow);}
</style>
<script>
(function(){
  function sync(){
    try{
      var p=window.parent,doc=p.document;
      var els=[doc.body,doc.documentElement];
      for(var i=0;i<els.length;i++){
        var cs=p.getComputedStyle(els[i]);
        var bg=cs.backgroundColor;
        var m=bg.match(/rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)/);
        if(!m) continue;
        var alpha=m[4]===undefined?1:parseFloat(m[4]);
        if(alpha<0.1) continue;
        var lum=(0.299*+m[1]+0.587*+m[2]+0.114*+m[3])/255;
        document.documentElement.classList.toggle('dark',lum<0.5);
        document.documentElement.classList.toggle('light',lum>=0.5);
        document.body.style.backgroundColor=bg;
        return;
      }
    }catch(e){}
  }
  sync();setInterval(sync,800);
})();
</script>
</head><body>
<div class="hero">
  <div class="eyebrow">Agente 03</div>
  <div class="title">Generaci&#243;n de <span class="accent">presentaciones</span></div>
  <div class="desc">Detecta ecuaciones y tablas en tu Markdown y genera PDF estructurado e interactivos HTML con Chart.js.</div>
  <div class="workflow">
    <div class="step">
      <div class="num">1</div>
      <div><div class="lbl">Paso 1</div><div class="sdesc">Markdown</div></div>
    </div>
    <svg class="arrow" width="16" height="10" viewBox="0 0 20 12" fill="none">
      <path d="M0 6h16M12 2l4 4-4 4" stroke="currentColor" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"/>
    </svg>
    <div class="step">
      <div class="num">2</div>
      <div><div class="lbl">Paso 2</div><div class="sdesc">Detectar</div></div>
    </div>
    <svg class="arrow" width="16" height="10" viewBox="0 0 20 12" fill="none">
      <path d="M0 6h16M12 2l4 4-4 4" stroke="currentColor" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"/>
    </svg>
    <div class="step">
      <div class="num">3</div>
      <div><div class="lbl">Paso 3</div><div class="sdesc">Seleccionar</div></div>
    </div>
    <svg class="arrow" width="16" height="10" viewBox="0 0 20 12" fill="none">
      <path d="M0 6h16M12 2l4 4-4 4" stroke="currentColor" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"/>
    </svg>
    <div class="step">
      <div class="num">4</div>
      <div><div class="lbl">Paso 4</div><div class="sdesc">Generar</div></div>
    </div>
  </div>
</div>
</body></html>"""


# ---------------------------------------------------------------------------
# Tipo icons y labels
# ---------------------------------------------------------------------------

_TIPO_ICON = {
    "ecuacion": "∑",
    "relacion": "f(x)",
    "tabla": "▦",
}
_TIPO_LABEL = {
    "ecuacion": "Ecuación",
    "relacion": "Relación paramétrica",
    "tabla": "Tabla numérica",
}

_TEXTO_ORIGINAL_MAX = 8000
_DENSIDAD_NUMERICA_RE = re.compile(
    r"\d+\.?\d*\s*(?:MPa|GPa|μm|µm|mm|nm|°|N|kN|Pa)",
    re.IGNORECASE,
)


def _extraer_texto_pdf_inteligente(file_bytes: bytes) -> str | None:
    """Prioriza páginas con más contenido numérico + las 2 primeras."""
    import pdfplumber

    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        pages_data: list[tuple[int, int, str]] = []
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            density = len(_DENSIDAD_NUMERICA_RE.findall(text))
            pages_data.append((i, density, text))

    if not pages_data:
        return None

    selected: set[int] = set()
    for i in range(min(2, len(pages_data))):
        selected.add(i)

    by_density = sorted(pages_data, key=lambda item: item[1], reverse=True)
    for i, density, _ in by_density:
        if len(selected) >= 5:
            break
        if density > 0 or len(selected) < 3:
            selected.add(i)

    ordered_text = "\n\n".join(
        pages_data[i][2] for i in sorted(selected) if pages_data[i][2].strip()
    )
    if not ordered_text.strip():
        return None
    return ordered_text[:_TEXTO_ORIGINAL_MAX]


def _extraer_texto_original(file_bytes: bytes, filename: str) -> str | None:
    """Extrae texto de PDF o PPTX; retorna None si falla o no hay texto."""
    try:
        name_lower = filename.lower()
        if name_lower.endswith(".pdf"):
            return _extraer_texto_pdf_inteligente(file_bytes)

        if name_lower.endswith(".pptx"):
            from pptx import Presentation

            parts: list[str] = []
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text.strip())
            if not parts:
                return None
            return "\n\n".join(parts)[:_TEXTO_ORIGINAL_MAX]
    except Exception:  # noqa: BLE001
        return None
    return None


# ---------------------------------------------------------------------------
# App principal
# ---------------------------------------------------------------------------

def main() -> None:
    # ── Session state ─────────────────────────────────────────────────────
    if "elementos" not in st.session_state:
        st.session_state["elementos"] = None
    if "md_content" not in st.session_state:
        st.session_state["md_content"] = None
    if "md_hash" not in st.session_state:
        st.session_state["md_hash"] = ""
    if "pdf_bytes" not in st.session_state:
        st.session_state["pdf_bytes"] = None
    if "pdf_titulo" not in st.session_state:
        st.session_state["pdf_titulo"] = "material"
    if "html_bytes" not in st.session_state:
        st.session_state["html_bytes"] = None
    if "html_titulo" not in st.session_state:
        st.session_state["html_titulo"] = "material"
    if "texto_original" not in st.session_state:
        st.session_state["texto_original"] = None

    st.set_page_config(page_title="Agente Presentación", layout="wide")

    # ── CSS global ────────────────────────────────────────────────────────
    st.markdown(
        """
<style>
@import url('https://fonts.googleapis.com/css2?family=Playfair+Display:wght@400;500&family=DM+Sans:wght@400;500&display=swap');

[data-testid="stAppViewContainer"] > .main,
[data-testid="stMain"] {
    background-color: var(--background-color) !important;
}
section[data-testid="stMain"] > div {
    background-color: var(--background-color) !important;
}
[data-testid="stSidebar"] {
    background-color: var(--secondary-background-color) !important;
    border-right: 1px solid rgba(128,128,128,0.2) !important;
}
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] {
    background-color: var(--secondary-background-color) !important;
    border-radius: 10px !important;
    border: 1px solid rgba(128,128,128,0.2) !important;
}
[data-testid="stFileUploaderDropzone"] {
    border-radius: 10px !important;
    border: 1px solid rgba(128,128,128,0.2) !important;
    background-color: var(--secondary-background-color) !important;
}
.stButton > button {
    background-color: #185FA5 !important;
    color: white !important;
    border: none !important;
    border-radius: 12px !important;
    font-family: 'DM Sans', sans-serif !important;
    font-weight: 500 !important;
    letter-spacing: 0.01em;
}
.stButton > button:hover {
    background-color: #0C447C !important;
}
.stButton > button:disabled {
    background-color: rgba(128,128,128,0.2) !important;
    color: var(--text-color) !important;
    opacity: 0.5 !important;
}
.stDownloadButton > button {
    border-radius: 12px !important;
    border: 1px solid rgba(128,128,128,0.2) !important;
    color: #185FA5 !important;
    font-family: 'DM Sans', sans-serif !important;
    font-weight: 500 !important;
}
.stDownloadButton > button:hover {
    background-color: rgba(24,95,165,0.05) !important;
}
[data-testid="stExpander"] {
    background-color: var(--secondary-background-color) !important;
    border: 0.5px solid rgba(128,128,128,0.2) !important;
    border-radius: 10px !important;
}
</style>
""",
        unsafe_allow_html=True,
    )

    # ── Sidebar ───────────────────────────────────────────────────────────
    with st.sidebar:
        # Branding
        st.markdown(
            """
<div style="padding-bottom:20px; border-bottom:1px solid rgba(128,128,128,0.2); margin-bottom:8px;">
  <div style="font-family:'DM Sans',sans-serif; font-size:16px; font-weight:500;
       color:var(--text-color); letter-spacing:-0.2px; line-height:1.2;">
    Agente Presentaci&#243;n
  </div>
</div>
""",
            unsafe_allow_html=True,
        )

        # Steps
        st.markdown(
            """
<div style="display:flex; flex-direction:column; gap:12px; margin-bottom:4px; padding-top:6px;">
  <div style="display:flex; align-items:center; gap:10px;">
    <span style="display:inline-flex; align-items:center; justify-content:center;
          width:20px; height:20px; border-radius:50%;
          background:#E6F1FB; color:#185FA5;
          font-family:'DM Sans',sans-serif; font-size:10px; font-weight:500; flex-shrink:0;">1</span>
    <span style="font-family:'DM Sans',sans-serif; font-size:12px; color:var(--text-color); opacity:0.7;">Sube el Markdown</span>
  </div>
  <div style="display:flex; align-items:center; gap:10px;">
    <span style="display:inline-flex; align-items:center; justify-content:center;
          width:20px; height:20px; border-radius:50%;
          background:#E6F1FB; color:#185FA5;
          font-family:'DM Sans',sans-serif; font-size:10px; font-weight:500; flex-shrink:0;">2</span>
    <span style="font-family:'DM Sans',sans-serif; font-size:12px; color:var(--text-color); opacity:0.7;">Detecta los elementos</span>
  </div>
  <div style="display:flex; align-items:center; gap:10px;">
    <span style="display:inline-flex; align-items:center; justify-content:center;
          width:20px; height:20px; border-radius:50%;
          background:#E6F1FB; color:#185FA5;
          font-family:'DM Sans',sans-serif; font-size:10px; font-weight:500; flex-shrink:0;">3</span>
    <span style="font-family:'DM Sans',sans-serif; font-size:12px; color:var(--text-color); opacity:0.7;">Selecciona y exporta</span>
  </div>
  <div style="display:flex; align-items:center; gap:10px;">
    <span style="display:inline-flex; align-items:center; justify-content:center;
          width:20px; height:20px; border-radius:50%;
          background:#E6F1FB; color:#185FA5;
          font-family:'DM Sans',sans-serif; font-size:10px; font-weight:500; flex-shrink:0;">4</span>
    <span style="font-family:'DM Sans',sans-serif; font-size:12px; color:var(--text-color); opacity:0.7;">Descarga el material</span>
  </div>
</div>
""",
            unsafe_allow_html=True,
        )

        st.divider()

        # File uploader
        st.markdown(
            """<div style="display:flex; align-items:center; gap:10px; margin:0 0 8px 0;">
  <span style="display:inline-flex; align-items:center; justify-content:center;
        width:20px; height:20px; border-radius:50%;
        background:#E6F1FB; color:#185FA5;
        font-family:'DM Sans',sans-serif; font-size:10px; font-weight:500; flex-shrink:0;">1</span>
  <div>
    <div style="font-family:'DM Sans',sans-serif; font-size:11px; font-weight:500;
         color:var(--text-color); letter-spacing:0.06em; text-transform:uppercase; line-height:1;">
      Markdown del tema</div>
    <div style="font-family:'DM Sans',sans-serif; font-size:12px; color:var(--text-color); opacity:0.55; margin-top:3px;">
      Archivo .md del Agente Contenido</div>
  </div>
</div>""",
            unsafe_allow_html=True,
        )
        uploaded_file = st.file_uploader(
            "Markdown del tema (.md)",
            type=["md"],
            accept_multiple_files=False,
            key="md_uploader",
        )

        # Detect file change and reset state
        if uploaded_file is not None:
            file_hash = hashlib.md5(uploaded_file.getvalue()).hexdigest()
            if file_hash != st.session_state["md_hash"]:
                st.session_state["elementos"] = None
                st.session_state["pdf_bytes"] = None
                st.session_state["html_bytes"] = None
                st.session_state["md_hash"] = file_hash
                st.session_state["md_content"] = uploaded_file.getvalue().decode(
                    "utf-8", errors="replace"
                )
                # Infer PDF title from filename
                stem = uploaded_file.name.replace(".md", "")
                st.session_state["pdf_titulo"] = stem

        st.markdown(
            """<div style="display:flex; align-items:center; gap:10px; margin:16px 0 8px 0;">
  <span style="display:inline-flex; align-items:center; justify-content:center;
        width:20px; height:20px; border-radius:50%;
        background:#E6F1FB; color:#185FA5;
        font-family:'DM Sans',sans-serif; font-size:10px; font-weight:500; flex-shrink:0;">+</span>
  <div>
    <div style="font-family:'DM Sans',sans-serif; font-size:11px; font-weight:500;
         color:var(--text-color); letter-spacing:0.06em; text-transform:uppercase; line-height:1;">
      Material original del profesor (opcional)</div>
    <div style="font-family:'DM Sans',sans-serif; font-size:12px; color:var(--text-color); opacity:0.55; margin-top:3px;">
      PDF o PPTX del tema. Mejora los rangos de los sliders
      y la selección del tipo de visualización.</div>
  </div>
</div>""",
            unsafe_allow_html=True,
        )
        uploaded_original = st.file_uploader(
            "Material original del profesor (opcional)",
            type=["pdf", "pptx"],
            accept_multiple_files=False,
            key="uploader_original",
            label_visibility="collapsed",
        )
        if uploaded_original is not None:
            st.session_state["texto_original"] = _extraer_texto_original(
                uploaded_original.getvalue(),
                uploaded_original.name,
            )
        else:
            st.session_state["texto_original"] = None

        st.divider()
        analizar_advertencias = st.checkbox(
            "Analizar advertencias pedagógicas (consume más créditos)",
            value=False,
            help="Opcional. Llama a Sonnet por cada elemento detectado para señalar "
            "posibles limitaciones interactivas antes de generar el HTML.",
            key="analizar_advertencias",
        )
        puede_detectar = uploaded_file is not None
        st.button(
            "Detectar elementos",
            key="detectar_btn",
            disabled=not puede_detectar,
            use_container_width=True,
        )
        st.button(
            "Generar PDF completo",
            key="generar_pdf_sidebar_btn",
            disabled=not puede_detectar,
            use_container_width=True,
        )
        if st.session_state.get("pdf_bytes") is not None:
            titulo_dl = st.session_state.get("pdf_titulo", "material")
            st.download_button(
                label="⬇ Descargar PDF",
                data=st.session_state["pdf_bytes"],
                file_name=f"{titulo_dl}_presentacion.pdf",
                mime="application/pdf",
                use_container_width=True,
                key="dl_pdf_sidebar",
            )

    # ── Area principal ────────────────────────────────────────────────────
    components.html(_HERO_PRES_HTML, height=340, scrolling=False)

    # Detection trigger
    if st.session_state.get("detectar_btn") and st.session_state["md_content"]:
        with st.status("Analizando el documento...", expanded=True) as status:
            st.write("🔍 Detectando ecuaciones y tablas...")
            try:
                elementos = detectar_elementos(
                    st.session_state["md_content"],
                    analizar_advertencias=st.session_state.get(
                        "analizar_advertencias", False
                    ),
                )
                st.session_state["elementos"] = elementos
                st.session_state["pdf_bytes"] = None  # reset previous PDF
                n = len(elementos)
                label = (
                    f"✅ {n} elementos detectados"
                    if n > 0
                    else "⚠️ No se detectaron elementos"
                )
                status.update(label=label, state="complete")
            except Exception as exc:
                status.update(label="❌ Error en la detección", state="error")
                st.error(f"Error: {exc}")
        st.rerun()

    # PDF generation trigger (sidebar button, independent of detection/selection)
    if st.session_state.get("generar_pdf_sidebar_btn") and st.session_state["md_content"]:
        _generated_ok = False
        with st.status("Generando PDF...", expanded=True) as status:
            md = st.session_state["md_content"] or ""
            titulo = st.session_state.get("pdf_titulo", "Material docente")
            try:
                st.write("📄 Renderizando ecuaciones y tablas...")
                pdf_bytes = generar_pdf(md, titulo=titulo)
                st.session_state["pdf_bytes"] = pdf_bytes
                status.update(label="✅ PDF generado", state="complete")
                _generated_ok = True
            except Exception as exc:
                status.update(label="❌ Error al generar el PDF", state="error")
                st.error(f"Error: {exc}")
        if _generated_ok:
            st.rerun()

    # No file uploaded
    if uploaded_file is None:
        st.info(
            "Sube un archivo **.md** generado por el Agente Contenido "
            "y pulsa **Detectar elementos** para comenzar."
        )
        return

    # File uploaded but no detection yet
    if st.session_state["elementos"] is None:
        st.info(
            "Archivo cargado. Pulsa **Detectar elementos** en el panel "
            "izquierdo para analizar el documento."
        )
        return

    elementos: list[dict] = st.session_state["elementos"]

    if not elementos:
        st.warning(
            "No se detectaron ecuaciones LaTeX ni tablas numéricas en el documento. "
            "Comprueba que el archivo es un output del Agente Contenido y contiene "
            "expresiones matemáticas en formato `$$...$$` o `$...$`."
        )
        return

    # ── Lista de elementos con checkboxes ─────────────────────────────────
    st.markdown("### Elementos detectados")
    st.caption(
        f"{len(elementos)} elementos — selecciona los que quieres incluir en el export."
    )

    # Group by type for visual clarity
    tipos_presentes = sorted({el["tipo"] for el in elementos})
    seleccionados: list[int] = []

    col_check_all, col_check_none, _ = st.columns([1, 1, 5])
    with col_check_all:
        if st.button("Seleccionar todo", key="sel_all", use_container_width=True):
            for el in elementos:
                st.session_state[f"chk_{el['id']}"] = True
    with col_check_none:
        if st.button("Ninguno", key="sel_none", use_container_width=True):
            for el in elementos:
                st.session_state[f"chk_{el['id']}"] = False

    st.markdown('<div style="height:8px;"></div>', unsafe_allow_html=True)

    for tipo in tipos_presentes:
        tipo_elementos = [el for el in elementos if el["tipo"] == tipo]
        icon = _TIPO_ICON.get(tipo, "·")
        label = _TIPO_LABEL.get(tipo, tipo)

        st.markdown(
            f'<div style="font-family: DM Sans, sans-serif; font-size:11px; '
            f'font-weight:500; color:var(--text-color); opacity:0.55; '
            f'text-transform:uppercase; letter-spacing:0.1em; margin:14px 0 6px 0;">'
            f'{icon} {label} ({len(tipo_elementos)})</div>',
            unsafe_allow_html=True,
        )

        for el in tipo_elementos:
            chk_key = f"chk_{el['id']}"
            # Initialize checkbox state on first render
            if chk_key not in st.session_state:
                st.session_state[chk_key] = True

            col_chk, col_exp = st.columns([1, 12])
            with col_chk:
                checked = st.checkbox(
                    el["nombre"],
                    key=chk_key,
                    label_visibility="collapsed",
                )
            with col_exp:
                st.markdown(f"**{el['nombre']}**")
                if el.get("advertencia"):
                    razon_esc = (
                        el["advertencia"]
                        .replace("&", "&amp;")
                        .replace("<", "&lt;")
                        .replace(">", "&gt;")
                    )
                    st.markdown(
                        f'<div style="font-size:12px;color:#E67E22;'
                        f'margin-bottom:4px;line-height:1.4;">'
                        f"⚠ Este contenido puede tener baja interactividad: "
                        f"{razon_esc}</div>",
                        unsafe_allow_html=True,
                    )
                with st.expander("Ver expresión y contexto", expanded=False):
                    # expresion contains $$...$$ blocks — use markdown for rendering
                    st.markdown(el["expresion"])
                    if el.get("contexto"):
                        st.caption(el["contexto"][:300])

            if checked:
                seleccionados.append(el["id"])

    n_sel = len(seleccionados)

    # ── Botones de generacion ─────────────────────────────────────────────
    st.divider()
    st.markdown("### Generar HTML interactivo")

    hay_seleccion = n_sel > 0
    if not hay_seleccion:
        st.caption(
            "Selecciona al menos una sección para generar el HTML interactivo. "
            "El PDF completo se genera desde la barra lateral, sin selección."
        )

    col_html, _ = st.columns([1, 1])
    with col_html:
        if st.button(
            f"Generar HTML ({n_sel} elementos)" if hay_seleccion else "Generar HTML interactivo",
            key="generar_html_btn",
            disabled=not hay_seleccion,
            use_container_width=True,
        ):
            with st.status("Generando HTML interactivo...", expanded=True) as status:
                try:
                    st.write("⚙️ Generando bloques con Sonnet (puede tardar ~30 s)...")
                    titulo = st.session_state.get("pdf_titulo", "Material interactivo")
                    seleccionados_set = set(seleccionados)
                    elementos_sel = [el for el in elementos if el["id"] in seleccionados_set]
                    html_str = generar_html(
                        elementos_sel,
                        titulo,
                        verbose=True,
                        texto_original=st.session_state.get("texto_original"),
                    )
                    st.session_state["html_bytes"] = html_str.encode("utf-8")
                    st.session_state["html_titulo"] = titulo
                    status.update(label="✅ HTML generado", state="complete")
                except Exception as exc:
                    status.update(label="❌ Error al generar el HTML", state="error")
                    st.error(f"Error: {exc}")
            st.rerun()

    # ── Descargas ─────────────────────────────────────────────────────────
    hay_descarga = (
        st.session_state["pdf_bytes"] is not None
        or st.session_state.get("html_bytes") is not None
    )
    if hay_descarga:
        st.divider()
        dl_col_pdf, dl_col_html, _ = st.columns([1, 1, 3])

        if st.session_state["pdf_bytes"] is not None:
            titulo = st.session_state.get("pdf_titulo", "material")
            with dl_col_pdf:
                st.download_button(
                    label="⬇ Descargar PDF",
                    data=st.session_state["pdf_bytes"],
                    file_name=f"{titulo}_presentacion.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                )

        if st.session_state.get("html_bytes") is not None:
            titulo = st.session_state.get("html_titulo", "material")
            filename = f"{_slug_html(titulo)}_interactivo.html"
            with dl_col_html:
                st.download_button(
                    label="⬇ Descargar HTML interactivo",
                    data=st.session_state["html_bytes"],
                    file_name=filename,
                    mime="text/html",
                    use_container_width=True,
                )


if __name__ == "__main__":
    main()
