"""Extraccion de texto desde PDF/PPTX."""

from __future__ import annotations

import logging
import re
import sys
from pathlib import Path

from cleaner import clean_extracted_text

_MONOREPO_ROOT = Path(__file__).resolve().parent.parent
if str(_MONOREPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_MONOREPO_ROOT))

_LOGGER = logging.getLogger(__name__)


def _ensure_extractor_audit_handler() -> None:
    if _LOGGER.handlers:
        return
    _LOGGER.setLevel(logging.INFO)
    handler = logging.StreamHandler()
    handler.setFormatter(logging.Formatter("%(levelname)s %(name)s: %(message)s"))
    _LOGGER.addHandler(handler)
    _LOGGER.propagate = False


_ensure_extractor_audit_handler()


_CONSONANT_RUN_RE = re.compile(r"(?i)[bcdfghjklmnñpqrstvwxyz]{5,}")
_KNOWN_WORDS = {
    "de",
    "del",
    "la",
    "el",
    "tema",
    "chapter",
    "section",
    "slide",
    "content",
    "engineering",
    "machine",
    "elements",
}


def _is_mirrored_text(line: str) -> bool:
    """
    Detecta texto extraído en espejo (rotado 180°).
    Heurística: una línea es texto en espejo si al invertirla
    forma palabras reales en español o inglés más que el original.
    Señales simples: secuencias de consonantes sin vocales >4 chars,
    o palabras conocidas al revertir.
    """
    stripped = line.strip()
    if len(stripped) < 4:
        return False
    reversed_line = stripped[::-1]
    vowels = set("aeiouáéíóúAEIOUÁÉÍÓÚ")
    original_vowels = sum(1 for c in stripped if c in vowels)
    reversed_vowels = sum(1 for c in reversed_line if c in vowels)
    if len(stripped) > 0:
        ratio = original_vowels / len(stripped)
        if ratio < 0.1 and reversed_vowels > original_vowels:
            return True

    original_tokens = [t.lower() for t in stripped.split()]
    reversed_tokens = [t.lower() for t in reversed_line.split()]
    original_known = sum(1 for t in original_tokens if t in _KNOWN_WORDS)
    reversed_known = sum(1 for t in reversed_tokens if t in _KNOWN_WORDS)
    if reversed_known > original_known and reversed_known > 0:
        return True

    if _CONSONANT_RUN_RE.search(stripped):
        rev_has_consonant_run = _CONSONANT_RUN_RE.search(reversed_line) is not None
        if not rev_has_consonant_run:
            return True
    return False


_PAGE_MARK_RE = re.compile(r"^\[PAGINA\s+(\d+)\]", re.MULTILINE | re.IGNORECASE)


def _clean_page_blocks(raw_document: str, filename: str, *, light: bool = False) -> str:
    """Aplica filtro de espejo y cleaner por bloque [PAGINA N]."""
    doc = (raw_document or "").strip()
    if not doc:
        return ""

    matches = list(_PAGE_MARK_RE.finditer(doc))
    if not matches:
        filtered = [ln for ln in doc.split("\n") if not _is_mirrored_text(ln)]
        return clean_extracted_text("\n".join(filtered), filename, light=light).strip()

    parts: list[str] = []
    for i, m in enumerate(matches):
        label = m.group(0).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(doc)
        body = doc[start:end].strip()
        if body == "[TEXTO_ILEGIBLE]":
            parts.append(f"{label}\n[TEXTO_ILEGIBLE]")
            continue
        filtered_lines = [ln for ln in body.split("\n") if not _is_mirrored_text(ln)]
        text = clean_extracted_text("\n".join(filtered_lines), filename, light=light).strip()
        if text:
            parts.append(f"{label}\n{text}")
        else:
            parts.append(f"{label}\n[TEXTO_ILEGIBLE]")
    return "\n\n".join(parts).strip()


def _extract_pdf_enriched(path: Path) -> str | None:
    """Extracción con jerarquía visual (#/##/###). Intenta pymupdf primero, luego pdfplumber."""
    # Intentar pymupdf: mejor decodificación de fuentes math y orden de lectura
    try:
        from shared.pdf_enriched import build_pdf_markdown_pymupdf

        raw = build_pdf_markdown_pymupdf(path)
        if raw:
            cleaned = _clean_page_blocks(raw, path.name, light=True)
            if cleaned.strip():
                _LOGGER.info("PDF pymupdf %s: %s chars", path.name, len(cleaned))
                return cleaned
    except Exception as exc:
        _LOGGER.warning("Extracción pymupdf falló, usando pdfplumber: %s", exc)

    # Fallback: pdfplumber enriquecido
    try:
        from shared.pdf_enriched import build_pdf_markdown

        raw = build_pdf_markdown(path)
        if not raw:
            return None
        cleaned = _clean_page_blocks(raw, path.name, light=True)
        return cleaned if cleaned.strip() else None
    except Exception as exc:
        _LOGGER.warning("Extracción PDF enriquecida no disponible: %s", exc)
        return None


def _extract_pdf_plain(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            try:
                texto_pagina = page.extract_text() or ""
                _LOGGER.info(
                    "Página %s: %s chars extraídos; primeros 200: %s",
                    idx,
                    len(texto_pagina),
                    repr(texto_pagina[:200]),
                )
                filtered_lines = [
                    ln for ln in texto_pagina.split("\n") if not _is_mirrored_text(ln)
                ]
                text = clean_extracted_text("\n".join(filtered_lines), path.name).strip()
                if text:
                    parts.append(f"[PAGINA {idx}]\n{text}")
                else:
                    _LOGGER.warning("Página %s: texto ilegible tras extracción/limpieza", idx)
                    parts.append(f"[PAGINA {idx}]\n[TEXTO_ILEGIBLE]")
            except Exception:
                _LOGGER.warning("Página %s: error en extracción, marcando ilegible", idx)
                parts.append(f"[PAGINA {idx}]\n[TEXTO_ILEGIBLE]")
    return "\n\n".join(parts).strip()


def _extract_pdf(path: Path) -> str:
    enriched = _extract_pdf_enriched(path)
    if enriched:
        n_headings = sum(
            1
            for ln in enriched.splitlines()
            if ln.lstrip().startswith(("# ", "## ", "### "))
        )
        _LOGGER.info(
            "PDF enriquecido %s: %s chars, %s líneas con prefijo de título",
            path.name,
            len(enriched),
            n_headings,
        )
        return enriched

    _LOGGER.info(
        "PDF %s: extracción plana (metadatos de fuente insuficientes)",
        path.name,
    )
    return _extract_pdf_plain(path)


def _escape_markdown_table_cell(text: str) -> str:
    """
    Escapa contenido de celda para tablas Markdown con pipes.
    Sin esto, un '|' literal en la celda rompe la tabla y confunde al pipeline.
    """
    return text.replace("\\", "\\\\").replace("|", "\\|")


def _pptx_table_to_markdown(table: object) -> str:
    """
    Convierte una tabla de python-pptx a Markdown tipo GitHub (pipes).
    La primera fila se trata como cabecera: es la convención más habitual en PPTX
    y coincide con el enfoque de MarkItDown (primera fila = th).
    """
    rows = list(getattr(table, "rows", []) or [])
    if not rows:
        return ""
    lines: list[str] = []
    for row_idx, row in enumerate(rows):
        cells = [
            _escape_markdown_table_cell(
                str(getattr(cell, "text", "") or "").strip().replace("\n", " ")
            )
            for cell in row.cells
        ]
        lines.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(lines)


def _pptx_collect_shape_blocks(shape: object, title_shape: object | None) -> list[str]:
    """
    Recorre shapes en el orden que expone python-pptx (orden Z de la diapositiva).
    Los grupos se expanden: el texto en grupos anidados no aparece en el iterador
    plano de slide.shapes y se perdería sin recursión.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

    blocks: list[str] = []

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in getattr(shape, "shapes", []):
            blocks.extend(_pptx_collect_shape_blocks(child, title_shape))
        return blocks

    # El título de slide ya va en un bloque `#` aparte: no duplicar el mismo shape.
    if title_shape is not None and shape is title_shape:
        return blocks

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE and getattr(
        shape, "table", None
    ):
        md_table = _pptx_table_to_markdown(shape.table)
        if md_table.strip():
            blocks.append(md_table.strip())
        return blocks

    if getattr(shape, "has_text_frame", False):
        raw = str(getattr(shape, "text", "") or "").strip()
        if not raw:
            return blocks
        # Subtítulo de layout típico: segunda jerarquía sin inferir roles complejos
        # (solo placeholders explícitos de PPTX, evitamos heurísticas frágiles).
        prefix = ""
        try:
            if getattr(shape, "is_placeholder", False):
                ph_type = shape.placeholder_format.type
                subtitle_types = {PP_PLACEHOLDER.SUBTITLE}
                if hasattr(PP_PLACEHOLDER, "VERTICAL_SUBTITLE"):
                    subtitle_types.add(PP_PLACEHOLDER.VERTICAL_SUBTITLE)
                if ph_type in subtitle_types:
                    prefix = "## "
        except (ValueError, AttributeError) as exc:
            _LOGGER.debug("PPTX placeholder no legible en slide: %s", exc)
        blocks.append(prefix + raw)
        return blocks

    return blocks


def _pptx_slide_to_text(slide: object) -> str:
    """
    Ensambla el contenido de una diapositiva: título como H1, cuerpo/tablas en orden,
    notas al final en un bloque propio (trazabilidad docente / guion del profesor).
    """
    blocks: list[str] = []

    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "text", None):
        title_text = str(title_shape.text).strip()
        if title_text:
            # Un solo H1 por slide: refleja el modelo mental de "título de diapositiva"
            # y alinea el Markdown con la jerarquía visual habitual.
            blocks.append(f"# {title_text}")

    for shape in slide.shapes:
        blocks.extend(_pptx_collect_shape_blocks(shape, title_shape))

    if getattr(slide, "has_notes_slide", False):
        notes_tf = slide.notes_slide.notes_text_frame
        if notes_tf is not None:
            notes_text = str(notes_tf.text or "").strip()
            if notes_text:
                # Bloque separado y estable para chunking/LLM sin mezclarlo con el bullet principal.
                blocks.append("### Notas del presentador\n\n" + notes_text)

    # Doble salto entre bloques: separa tablas Markdown del texto plano sin fusionar filas.
    return "\n\n".join(b for b in blocks if b).strip()


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        try:
            slide_raw_text = _pptx_slide_to_text(slide)
            if slide_raw_text:
                _LOGGER.info(
                    "Slide %s: %s chars extraídos; primeros 200: %s",
                    idx,
                    len(slide_raw_text),
                    repr(slide_raw_text[:200]),
                )
                filtered_lines = [
                    ln for ln in slide_raw_text.split("\n") if not _is_mirrored_text(ln)
                ]
                slide_text = clean_extracted_text("\n".join(filtered_lines), path.name).strip()
                if slide_text:
                    parts.append(f"[SLIDE {idx}]\n{slide_text}")
                else:
                    _LOGGER.warning("Slide %s: texto ilegible tras extracción/limpieza", idx)
                    parts.append(f"[SLIDE {idx}]\n[TEXTO_ILEGIBLE]")
            else:
                _LOGGER.warning("Slide %s: sin texto extraíble", idx)
                parts.append(f"[SLIDE {idx}]\n[TEXTO_ILEGIBLE]")
        except Exception:
            _LOGGER.warning("Slide %s: error en extracción, marcando ilegible", idx)
            parts.append(f"[SLIDE {idx}]\n[TEXTO_ILEGIBLE]")
    return "\n\n".join(parts).strip()


def extract_text(file_path: str) -> str:
    """Extrae texto de un PDF o PPTX."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"No existe el archivo: {file_path}")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".pptx":
        return _extract_pptx(path)

    raise ValueError("Formato no soportado. Usa PDF o PPTX.")
