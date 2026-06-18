import difflib
import io
import re
import unicodedata
from collections import Counter, defaultdict
from pathlib import Path

import pdfplumber
from pptx import Presentation


def clasificar_archivo(nombre: str, texto: str) -> str:
    nombre_normalizado = (nombre or "").lower()
    muestra_texto = (texto or "")[:300].lower()

    indicadores_nombre = [
        "outline",
        "index",
        "program",
        "syllabus",
        "c0",
        "tema0",
        "indice",
    ]
    indicadores_texto = [
        "outline",
        "contents",
        "index",
        "programa",
        "índice general",
        "tabla de contenidos",
    ]

    if any(indicador in nombre_normalizado for indicador in indicadores_nombre):
        return "contexto"
    if any(indicador in muestra_texto for indicador in indicadores_texto):
        return "contexto"

    return "teoria"


def extraer_texto(archivo_bytes, nombre_archivo) -> str:
    extension = Path(nombre_archivo).suffix.lower()

    if extension == ".pdf":
        texto_paginas = []
        with pdfplumber.open(io.BytesIO(archivo_bytes)) as pdf:
            for indice, pagina in enumerate(pdf.pages, start=1):
                texto = (pagina.extract_text() or "").strip()
                texto_paginas.append(f"--- Página {indice} ---\n{texto}")

        texto_final = "\n\n".join(texto_paginas).strip()
    elif extension == ".pptx":
        presentacion = Presentation(io.BytesIO(archivo_bytes))
        texto_slides = []

        for indice, slide in enumerate(presentacion.slides, start=1):
            fragmentos = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto_shape = (shape.text or "").strip()
                    if texto_shape:
                        fragmentos.append(texto_shape)

            texto_slide = "\n".join(fragmentos).strip()
            texto_slides.append(f"--- Slide {indice} ---\n{texto_slide}")

        texto_final = "\n\n".join(texto_slides).strip()
    else:
        raise ValueError(f"Formato no soportado para '{nombre_archivo}'. Solo se admite PDF o PPTX.")

    if not texto_final or len(texto_final) < 50:
        raise ValueError(
            f"No se pudo extraer texto suficiente de '{nombre_archivo}'. "
            "El archivo está vacío, protegido o no contiene texto legible."
        )

    return texto_final


# ---------------------------------------------------------------------------
# Detección determinista de subtemas por numeración jerárquica.
# Patrón replicado de agente-contenido/chunker.py (_NUMBERED_HEADER_RE).
# ---------------------------------------------------------------------------

_SUBTEMA_NUM_RE = re.compile(r"^\d+(?:\.\d+)*\.?\s+\S")
_PAGINA_FOOTER_RE = re.compile(r"^\d+\s+de\s+\d+$", re.I)


def normalizar_subtema(texto: str) -> str:
    """Elimina prefijo numérico, acentos y puntuación; convierte a minúsculas.
    Usada para comparar candidatos entre fuentes (guía vs material).
    """
    s = re.sub(r"^\d+(?:\.\d+)*\.?\s*", "", (texto or "")).lower()
    s = unicodedata.normalize("NFD", s)
    s = "".join(c for c in s if unicodedata.category(c) != "Mn")
    return re.sub(r"[^a-z0-9 ]", " ", s).strip()


# ---------------------------------------------------------------------------
# Filtro de calidad de subtemas candidatos.
#
# Causa raíz de dos fallos observados (Elementos de Máquinas, junio 2026):
#  1. Boilerplate administrativo de las guías docentes UniOvi tratado como
#     subtema técnico (Identificación, Contextualización, Requisitos,
#     Competencias, Metodología, Evaluación, Recursos/bibliografía). Estas
#     secciones se numeran en la guía y el regex de numeración las capturaba.
#  2. Fragmentos de texto corrido del material (p. ej. un paso de un ejemplo
#     resuelto "5.6. Además, de acuerdo a la tabla...") aceptados como subtema
#     porque empezaban por un número, sin validar que fueran un encabezado.
#
# Decisión de diseño documentada: un subtema candidato debe ser un encabezado
# real, no boilerplate administrativo ni prosa. Estos nombres administrativos
# son muy estables entre guías de la misma universidad, así que se filtran por
# patrón; la prosa se descarta por señales estructurales (conector inicial,
# ruido numérico, fragmento truncado).
# ---------------------------------------------------------------------------

# Secciones administrativas estándar de las guías docentes UniOvi. Prefijos
# (distintivos, no prefijan contenido técnico real) y nombres exactos cortos.
_GUIA_SECCIONES_ADMIN_PREFIJOS = (
    "identificacion de la asignatura",
    "competencias y resultados de aprendizaje",
    "metodologia y plan de trabajo",
    "evaluacion del aprendizaje",
    "recursos bibliografia",
    "distribucion de los contenidos",
)
_GUIA_SECCIONES_ADMIN_EXACTAS = {
    "contenidos", "contextualizacion", "requisitos", "competencias",
    "metodologia", "evaluacion", "recursos", "bibliografia",
    "presenciales", "no presenciales", "clases expositivas",
    "practicas de aula", "practicas de aula seminarios",
    "practicas de laboratorio", "practicas de laboratorio campo",
    "tutorias grupales", "sesiones de evaluacion",
    "trabajo autonomo", "trabajo en grupo",
    "exposicion de trabajos realizados en grupo",
    "evaluacion continua", "evaluacion extraordinaria", "evaluacion diferenciada",
}

# Conectores discursivos: una frase que empieza por uno de estos es prosa
# (texto corrido), nunca un título/encabezado de subtema.
_CONECTORES_PROSA = (
    "ademas", "asimismo", "por tanto", "por lo tanto", "sin embargo",
    "no obstante", "es decir", "de acuerdo", "por ejemplo", "en consecuencia",
    "por consiguiente", "entonces", "asi pues", "de hecho", "en este caso",
    "con ello", "con lo que", "de esta forma", "de esta manera", "de este modo",
    "por ultimo", "finalmente", "en resumen", "en definitiva", "a continuacion",
)

# Palabras-función con las que un encabezado real no termina; si una línea corta
# acaba en una de ellas es un fragmento truncado (típico de columnas de tabla).
_PALABRAS_FUNCION_FINAL = {
    "de", "del", "la", "el", "los", "las", "y", "o", "en", "a", "con",
    "por", "para", "que", "su", "sus", "un", "una", "al",
}


def _es_seccion_administrativa(norm: str) -> bool:
    """True si `norm` (nombre ya normalizado) es una sección administrativa de guía."""
    if norm in _GUIA_SECCIONES_ADMIN_EXACTAS:
        return True
    return any(norm.startswith(p) for p in _GUIA_SECCIONES_ADMIN_PREFIJOS)


def es_subtema_valido(nombre: str) -> bool:
    """True si `nombre` parece un subtema/encabezado técnico real.

    Filtro de calidad común a todos los caminos de detección. Rechaza:
      1. Secciones administrativas estándar de la guía docente UniOvi.
      2. Fragmentos de prosa que empiezan por un conector discursivo.
      3. Filas de datos numéricos (p. ej. filas de tablas de horas).
      4. Fragmentos cortos y truncados que terminan en palabra-función.
    """
    norm = normalizar_subtema(nombre)
    if not norm:
        return False
    if _es_seccion_administrativa(norm):
        return False
    if any(norm == c or norm.startswith(c + " ") for c in _CONECTORES_PROSA):
        return False
    tokens = norm.split()
    # Encabezados demasiado cortos o tokens sueltos (p. ej. «L», «R», «N»).
    if len(norm) < 4:
        return False
    if len(tokens) == 1 and len(tokens[0]) <= 2:
        return False
    # Fragmentos numéricos / unidades (p. ej. «l/100», «0.2 l/100» → «l/100»).
    compacto = norm.replace(" ", "")
    if re.fullmatch(r"[\dl./\\-]+", compacto):
        return False
    if "/" in nombre and len(norm) < 12:
        return False
    # Ruido numérico: 3+ tokens con dígitos → fila de tabla / fragmento de cálculo.
    n_num = sum(1 for t in tokens if any(ch.isdigit() for ch in t))
    if n_num >= 3:
        return False
    # Fragmento truncado: corto y acaba en palabra-función.
    if len(norm) <= 30 and tokens and tokens[-1] in _PALABRAS_FUNCION_FINAL:
        return False
    return True


def extraer_subtemas_candidatos(texto: str) -> list[str]:
    """Detecta subtemas con numeración jerárquica ('3.2. Título').
    Devuelve nombres limpios (sin prefijo numérico), deduplicados por normalización.
    Aplica el filtro de calidad es_subtema_valido() para descartar boilerplate
    administrativo y fragmentos de texto corrido.
    """
    candidatos: list[str] = []
    vistos: set[str] = set()
    for linea in (texto or "").splitlines():
        linea = linea.strip()
        if _PAGINA_FOOTER_RE.match(linea):
            continue
        if not _SUBTEMA_NUM_RE.match(linea):
            continue
        nombre = re.sub(r"^\d+(?:\.\d+)*\.?\s*", "", linea).strip()
        if not nombre or not es_subtema_valido(nombre):
            continue
        clave = normalizar_subtema(nombre)
        if clave and clave not in vistos:
            candidatos.append(nombre)
            vistos.add(clave)
    return candidatos


def extraer_subtemas_guia(texto: str) -> list[str]:
    """Extrae los subtemas reales de la sección 'Contenidos' de una guía docente.

    Las guías UniOvi numeran TODAS sus secciones de primer nivel (1..8), de las
    cuales solo 'Contenidos' contiene temario real; el resto es boilerplate
    administrativo. Acotar la extracción a la sección 'Contenidos' excluye de
    raíz cabeceras, filas de tablas de horas y prosa de metodología/evaluación.

    Estrategia: localizar la cabecera 'N. Contenidos' y recoger los subtemas
    numerados hasta la siguiente cabecera de sección administrativa. Si no se
    encuentra 'Contenidos' (guía con otro formato), se degrada a
    extraer_subtemas_candidatos() sobre todo el texto. En ambos caminos se aplica
    el filtro de calidad es_subtema_valido().
    """
    lineas = [l.strip() for l in (texto or "").splitlines()]

    inicio = None
    for i, ln in enumerate(lineas):
        m = re.match(r"^\d+\.\s+(\S.*)$", ln)
        if m and normalizar_subtema(m.group(1)).startswith("contenidos"):
            inicio = i
            break

    if inicio is None:
        # Formato no reconocido: degradar a extracción genérica (ya filtrada).
        return extraer_subtemas_candidatos(texto)

    candidatos: list[str] = []
    vistos: set[str] = set()
    for ln in lineas[inicio + 1:]:
        if _PAGINA_FOOTER_RE.match(ln):
            continue
        if not _SUBTEMA_NUM_RE.match(ln):
            continue
        nombre = re.sub(r"^\d+(?:\.\d+)*\.?\s*", "", ln).strip()
        norm = normalizar_subtema(nombre)
        if not norm:
            continue
        # Fin de la sección 'Contenidos': siguiente cabecera administrativa.
        if _es_seccion_administrativa(norm):
            break
        if not es_subtema_valido(nombre) or norm in vistos:
            continue
        candidatos.append(nombre)
        vistos.add(norm)
    return candidatos


def hay_discrepancia(lista_a: list[str], lista_b: list[str]) -> bool:
    """True si A y B no comparten ningún elemento (normalizado).
    Señal de revisión cuando la guía y el material describen el bloque de forma distinta.
    """
    if not lista_a or not lista_b:
        return False
    norm_a = {normalizar_subtema(x) for x in lista_a if x}
    norm_b = {normalizar_subtema(x) for x in lista_b if x}
    return norm_a.isdisjoint(norm_b)


# ---------------------------------------------------------------------------
# Extracción de señales estructurales con evidencia verificable.
# Objetivo: cada subbloque generado puede justificarse señalando una referencia
# concreta del documento fuente (número de sección, slide, página).
# ---------------------------------------------------------------------------


def extraer_titulos_slides_pptx(archivo_bytes: bytes) -> list[dict]:
    """Extrae títulos de diapositiva de un PPTX como señales estructurales.

    Prioridad por shape:
    1. Placeholder con idx=0 (marcador oficial de título en la plantilla).
    2. Primer texto corto sin saltos de línea (≤ 100 caracteres) — heurística
       para slides sin placeholder de título explícito.

    Los títulos duplicados (por normalización) se descartan para no generar
    candidatos repetidos cuando varias slides comparten encabezado.

    Retorna [{slide, titulo, referencia}].
    """
    presentacion = Presentation(io.BytesIO(archivo_bytes))
    resultados: list[dict] = []
    vistos: set[str] = set()

    for indice, slide in enumerate(presentacion.slides, start=1):
        titulo_oficial: str | None = None
        primer_texto_corto: str | None = None

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            texto = (shape.text or "").strip()
            if not texto:
                continue

            ph = getattr(shape, "placeholder_format", None)
            if ph is not None and ph.idx == 0:
                titulo_oficial = texto
                break

            if primer_texto_corto is None and len(texto) <= 100 and "\n" not in texto:
                primer_texto_corto = texto

        titulo_final = titulo_oficial or primer_texto_corto
        if not titulo_final:
            continue

        clave = normalizar_subtema(titulo_final)
        if not clave or clave in vistos:
            continue
        vistos.add(clave)

        resultados.append({
            "slide": indice,
            "titulo": titulo_final,
            "referencia": f"Slide {indice}",
        })

    return resultados


def extraer_titulos_visuales_pdf(archivo_bytes: bytes) -> list[dict]:
    """Detecta líneas candidatas a título en un PDF por tamaño/estilo de fuente.

    Heurística de frecuencia (SciPlore Xtract): la combinación (fontname, size)
    más frecuente en el documento es el cuerpo; las líneas cortas y uniformes
    cuya combinación tiene mayor tamaño o nombre de fuente negrita son candidatas
    a título. Esta heurística simple supera a un SVM entrenado para la misma
    tarea en el benchmark de referencia (77,9 % frente a 69,4 %).

    Notas sobre pdfplumber:
    - Los tamaños son relativos al documento (pueden diferir de Adobe Acrobat por
      un offset conocido, irrelevante para comparaciones internas).
    - La negrita se infiere de substrings "bold"/"bd"/"black"/"heavy" en fontname;
      pdfplumber no expone un booleano de negrita.

    Retorna [{titulo: str, referencia: str}] con referencia = "p. N".
    """
    resultados: list[dict] = []
    vistos_norm: set[str] = set()

    try:
        with pdfplumber.open(io.BytesIO(archivo_bytes)) as pdf:
            # 1. Recoger todas las palabras con metadatos de fuente
            todas_palabras: list[dict] = []
            alturas_pagina: dict[int, float] = {}

            for npag, pagina in enumerate(pdf.pages, start=1):
                try:
                    palabras = pagina.extract_words(extra_attrs=["fontname", "size"]) or []
                except Exception:
                    palabras = []
                alturas_pagina[npag] = float(pagina.height or 800)
                for p in palabras:
                    p["_pagina"] = npag
                    todas_palabras.append(p)

            # PDF escaneado o vacío: no hay metadatos de fuente útiles
            if len(todas_palabras) < 20:
                return []

            # 2. Determinar (fontname, size) del cuerpo por frecuencia
            conteo: Counter = Counter()
            for p in todas_palabras:
                fn = (p.get("fontname") or "").strip()
                sz = round(float(p.get("size") or 0) * 2) / 2  # resolución 0.5 pt
                if fn and sz > 0:
                    conteo[(fn, sz)] += 1

            if not conteo:
                return []

            cuerpo_fn, cuerpo_sz = conteo.most_common(1)[0][0]

            # 3. Agrupar palabras en líneas por (página, y redondeado)
            lineas: dict[tuple, list[dict]] = defaultdict(list)
            for p in todas_palabras:
                y_bucket = int(round(float(p.get("top") or 0)))
                lineas[(p["_pagina"], y_bucket)].append(p)

            # 4. Evaluar cada línea como candidata a título
            for (npag, y_bucket), palabras_linea in sorted(lineas.items()):
                palabras_linea.sort(key=lambda w: float(w.get("x0") or 0))
                texto = " ".join(
                    w["text"] for w in palabras_linea if w.get("text")
                ).strip()

                # Descartar líneas demasiado largas para ser un título
                if not texto or len(texto) > 120 or len(texto.split()) > 15:
                    continue

                # Excluir márgenes: cabeceras y pies de página (~6 % arriba/abajo)
                altura_pag = alturas_pagina.get(npag, 800.0)
                margen = altura_pag * 0.06
                if y_bucket < margen or y_bucket > altura_pag - margen:
                    continue

                # Estilo dominante de la línea (fontname, size más frecuente)
                estilos: Counter = Counter()
                for w in palabras_linea:
                    fn = (w.get("fontname") or "").strip()
                    sz = round(float(w.get("size") or 0) * 2) / 2
                    if fn and sz > 0:
                        estilos[(fn, sz)] += 1

                if not estilos:
                    continue

                (fn_dom, sz_dom), n_dom = estilos.most_common(1)[0]
                n_total = sum(estilos.values())

                # Línea con estilo mixto (≥ 30 % palabras fuera del estilo dominante) → prosa
                if n_dom < n_total * 0.7:
                    continue

                # El estilo dominante es el cuerpo → no es título
                if fn_dom == cuerpo_fn and sz_dom == cuerpo_sz:
                    continue

                # Criterio 1: tamaño mayor al cuerpo
                es_mayor = sz_dom > cuerpo_sz + 0.5
                # Criterio 2: nombre de fuente con indicador de peso (negrita/heavy)
                fn_dom_lower = fn_dom.lower()
                es_negrita = fn_dom != cuerpo_fn and any(
                    sub in fn_dom_lower for sub in ("bold", "bd", "black", "heavy")
                )

                if not (es_mayor or es_negrita):
                    continue

                if not es_subtema_valido(texto):
                    continue

                clave = normalizar_subtema(texto)
                if not clave or clave in vistos_norm:
                    continue

                vistos_norm.add(clave)
                resultados.append({"titulo": texto, "referencia": f"p. {npag}"})

    except Exception:
        return []

    return resultados


def extraer_candidatos_con_evidencia(
    texto: str,
    nombre_archivo: str = "",
    archivo_bytes: bytes | None = None,
) -> list[dict]:
    """Detecta subbloques con referencia estructural verificable.

    Prioridad de fuentes (estricta — si la primera produce resultados, no se
    consulta la siguiente):

    1. Secciones numeradas en el texto extraído (e.g. '3.2. Título').
       Aplica a PDF y PPTX. Evidencia: 'Sección 3.2'.
    2. Títulos de diapositiva en PPTX (placeholder idx=0 o heurística).
       Solo si no se encontraron secciones numeradas. Evidencia: 'Slide N'.
    3. Títulos visuales por tamaño/estilo de fuente en PDF (heurística de
       frecuencia). Solo para PDF sin secciones numeradas. Evidencia: 'p. N'.

    Retorna [] si ninguna fuente ofrece señales verificables — el bloque debe
    tratarse como un único subbloque (fallback); el caller es responsable de
    marcarlo en la interfaz.

    Cada ítem: {nombre: str, evidencia: str, fuente: str}
      - fuente: 'numeracion' | 'titulo_slide' | 'titulo_visual'
    """
    candidatos: list[dict] = []
    vistos: set[str] = set()

    # Estrategia 1 — secciones numeradas
    for linea in (texto or "").splitlines():
        linea = linea.strip()
        if _PAGINA_FOOTER_RE.match(linea):
            continue
        if not _SUBTEMA_NUM_RE.match(linea):
            continue
        m_pref = re.match(r"^(\d+(?:\.\d+)*)\.?\s*", linea)
        prefijo = m_pref.group(1) if m_pref else ""
        nombre = re.sub(r"^\d+(?:\.\d+)*\.?\s*", "", linea).strip()
        if not nombre or not es_subtema_valido(nombre):
            continue
        clave = normalizar_subtema(nombre)
        if not clave or clave in vistos:
            continue
        evidencia = f"Sección {prefijo}" if prefijo else "Sección numerada"
        candidatos.append({"nombre": nombre, "evidencia": evidencia, "fuente": "numeracion"})
        vistos.add(clave)

    if candidatos:
        return candidatos

    ext = Path(nombre_archivo).suffix.lower() if nombre_archivo else ""

    # Estrategia 2 — títulos de diapositiva PPTX
    if ext == ".pptx" and archivo_bytes is not None:
        for item in extraer_titulos_slides_pptx(archivo_bytes):
            if not es_subtema_valido(item["titulo"]):
                continue
            clave = normalizar_subtema(item["titulo"])
            if not clave or clave in vistos:
                continue
            candidatos.append({
                "nombre": item["titulo"],
                "evidencia": item["referencia"],
                "fuente": "titulo_slide",
            })
            vistos.add(clave)

    if candidatos:
        return candidatos

    # Estrategia 3 — títulos visuales por tamaño/estilo de fuente (PDF sin numeración)
    if ext == ".pdf" and archivo_bytes is not None:
        for item in extraer_titulos_visuales_pdf(archivo_bytes):
            clave = normalizar_subtema(item["titulo"])
            if not clave or clave in vistos:
                continue
            candidatos.append({
                "nombre": item["titulo"],
                "evidencia": item["referencia"],
                "fuente": "titulo_visual",
            })
            vistos.add(clave)

    return candidatos  # lista vacía → fallback obligatorio


_EVIDENCIA_VACIA = frozenset({
    "",
    "—",
    "-",
    "–",
    "sin señal verificable",
    "sin senal verificable",
    "fallback",
    "sin señal",
    "sin senal",
    "guía docente",
    "guia docente",
})


def evidencia_es_vacia(evidencia: str) -> bool:
    """True si la evidencia no aporta una frontera estructural utilizable."""
    return (evidencia or "").strip().lower() in _EVIDENCIA_VACIA


def _mapa_evidencia_candidatos(candidatos: list[dict]) -> dict[str, str]:
    """Índice nombre_normalizado → evidencia para un material."""
    mapa: dict[str, str] = {}
    for cand in candidatos:
        clave = normalizar_subtema(cand.get("nombre", ""))
        ev = (cand.get("evidencia") or "").strip()
        if clave and ev and not evidencia_es_vacia(ev):
            mapa[clave] = ev
    return mapa


# ---------------------------------------------------------------------------
# Subtemas confirmados para el prompt (app-unificada / standalone)
# ---------------------------------------------------------------------------


def construir_subtemas_confirmados(
    texto_guia: str,
    textos_teoria: list[str],
    archivos_teoria: list[str],
    archivos_bytes: list[bytes],
) -> list[list[dict]] | None:
    """Lista cerrada de subtemas por material para el prompt del Organizador.

    Si la guía docente aporta bloques temáticos (sección Contenidos), devuelve
    ``None`` para que el LLM estructure subtemas libremente desde guía + material
    (comportamiento del standalone validado). Solo impone lista cerrada cuando
    no hay guía y el material tiene señales estructurales fiables.
    """
    candidatos_guia = extraer_subtemas_guia(texto_guia)
    if len(candidatos_guia) >= 2:
        return None

    resultado: list[list[dict]] = []
    for texto_mat, nombre_arch, bytes_arch in zip(
        textos_teoria, archivos_teoria, archivos_bytes
    ):
        cands_mat = extraer_candidatos_con_evidencia(
            texto_mat, nombre_arch, bytes_arch
        )
        if cands_mat:
            resultado.append([
                {
                    "nombre": c["nombre"],
                    "evidencia": c["evidencia"],
                    "origen": "Detectado",
                }
                for c in cands_mat
            ])
        else:
            resultado.append([])

    return resultado


# ---------------------------------------------------------------------------
# Enriquecimiento de evidencia al confirmar la organización (app-unificada).
# La detección determinista ocurre antes del LLM; el Markdown generado suele
# dejar la columna Evidencia en «—». Estas funciones reasignan la evidencia
# real sin depender de que el modelo la copie fielmente.
# ---------------------------------------------------------------------------


def resolver_archivo_bloque(nombre_bloque: str, archivos: list[str], idx: int) -> str:
    """Asocia un bloque temático al PDF/PPTX de teoría más probable."""
    if not archivos:
        return ""
    nb = normalizar_subtema(nombre_bloque)
    mejor, score = "", 0.0
    for ar in archivos:
        stem_norm = normalizar_subtema(Path(ar).stem)
        r = difflib.SequenceMatcher(None, nb, stem_norm).ratio()
        if r > score:
            score, mejor = r, ar
    if score >= 0.35:
        return mejor
    return archivos[idx] if idx < len(archivos) else archivos[0]


def _buscar_evidencia_subtema(nombre: str, mapa: dict[str, str]) -> str:
    """Resuelve la evidencia de un subtema por coincidencia exacta de nombre."""
    clave = normalizar_subtema(nombre)
    if not clave:
        return ""
    return mapa.get(clave, "")


def enriquecer_bloques_con_evidencia_detectada(
    bloques: list[dict],
    candidatos_por_material: list[list[dict]],
    archivos_teoria: list[str],
) -> list[dict]:
    """Completa evidencia vacía en subtemas con la detección determinista del material.

    Añade ``archivo_origen`` a cada bloque para persistir el vínculo bloque→PDF.
  """
    if not bloques:
        return bloques

    mapas_por_archivo = {
        archivos_teoria[i]: _mapa_evidencia_candidatos(
            candidatos_por_material[i] if i < len(candidatos_por_material) else []
        )
        for i in range(len(archivos_teoria))
    }

    enriquecidos: list[dict] = []
    for idx_bloque, bloque in enumerate(bloques):
        bloque_out = dict(bloque)
        archivo = resolver_archivo_bloque(
            bloque_out.get("nombre", ""), archivos_teoria, idx_bloque
        )
        bloque_out["archivo_origen"] = archivo
        mapa = mapas_por_archivo.get(archivo, {})

        subtemas_out: list[dict] = []
        for sub in bloque_out.get("subtemas", []):
            sub_out = dict(sub)
            if evidencia_es_vacia(sub_out.get("evidencia", "")):
                ev = _buscar_evidencia_subtema(sub_out.get("nombre", ""), mapa)
                if ev:
                    sub_out["evidencia"] = ev
            subtemas_out.append(sub_out)
        bloque_out["subtemas"] = subtemas_out
        enriquecidos.append(bloque_out)

    return enriquecidos


# ---------------------------------------------------------------------------
# Serialización / deserialización de la organización para edición manual.
# ---------------------------------------------------------------------------


def _es_celda_horas(texto: str) -> bool:
    """True si el texto de una celda parece una hora (p. ej. '2.5h' o '3')."""
    raw = (texto or "").strip().lower().rstrip("h").strip()
    if not raw:
        return False
    try:
        float(raw.replace(",", "."))
        return True
    except ValueError:
        return False


def parsear_bloques_desde_markdown(markdown: str) -> list[dict]:
    """Parsea el Markdown de organización en una lista de bloques estructurados.

    Retorna [{numero, nombre, horas, subtemas: [{nombre, evidencia, origen, manual, aprobado}], manual}].

    Compatible con el formato actual (| Subtema | Evidencia | Origen |) y con formatos
    legados que incluían columna Horas por subtema.
    """
    BLOQUE_RE = re.compile(r"^## Bloque (\d+) — (.+?) · ([\d.]+)h", re.MULTILINE)
    _EXCLUIR = {
        "subtema", "topic", "horas", "hours",
        "justificación", "justificacion", "justification",
        "evidencia", "evidence", "origen", "origin",
    }

    bloques: list[dict] = []
    matches = list(BLOQUE_RE.finditer(markdown))

    for i, m in enumerate(matches):
        numero = int(m.group(1))
        nombre = m.group(2).strip()
        horas = float(m.group(3))
        inicio = m.end()
        fin = matches[i + 1].start() if i + 1 < len(matches) else len(markdown)
        seccion = markdown[inicio:fin]

        subtemas: list[dict] = []
        for linea in seccion.splitlines():
            linea = linea.strip()
            if not linea.startswith("|") or not linea.endswith("|"):
                continue
            celdas = _parsear_celdas(linea)
            if len(celdas) < 2:
                continue
            col1 = celdas[0].strip()
            if col1.lower() in _EXCLUIR or re.match(r"^-+$", col1.replace(" ", "")):
                continue
            if not col1:
                continue

            if len(celdas) >= 4:
                # Legado: | Subtema | Horas | Evidencia | Origen |
                subtemas.append({
                    "nombre": col1,
                    "evidencia": celdas[2].strip(),
                    "origen": celdas[3].strip(),
                    "manual": False,
                    "aprobado": False,
                })
            elif len(celdas) == 3 and _es_celda_horas(celdas[1]):
                # Legado: | Subtema | Horas | Origen |
                subtemas.append({
                    "nombre": col1,
                    "evidencia": "",
                    "origen": celdas[2].strip(),
                    "manual": False,
                    "aprobado": False,
                })
            elif len(celdas) >= 3:
                # Actual: | Subtema | Evidencia | Origen |
                subtemas.append({
                    "nombre": col1,
                    "evidencia": celdas[1].strip(),
                    "origen": celdas[2].strip(),
                    "manual": False,
                    "aprobado": False,
                })
            elif len(celdas) == 2:
                # Sin origen explícito: | Subtema | Evidencia |
                subtemas.append({
                    "nombre": col1,
                    "evidencia": celdas[1].strip(),
                    "origen": "Detectado",
                    "manual": False,
                    "aprobado": False,
                })

        bloques.append({
            "numero": numero,
            "nombre": nombre,
            "horas": horas,
            "subtemas": subtemas,
            "manual": False,
        })

    return bloques


def regenerar_markdown_desde_bloques(
    bloques: list[dict],
    markdown_original: str,
) -> str:
    """Regenera el Markdown de organización a partir del estado estructurado.

    Preserva la cabecera (título + resumen de horas) y el pie (nota PL) del
    markdown original. Las horas de cada bloque vienen del campo bloque['horas']
    (solo a nivel de bloque, no por subtema).

    Formato de salida: | Subtema | Evidencia | Origen |
    """
    BLOQUE_RE = re.compile(r"^## Bloque \d+", re.MULTILINE)
    FOOTER_RE = re.compile(r"^> ", re.MULTILINE)

    primer_m = BLOQUE_RE.search(markdown_original)
    header = markdown_original[: primer_m.start()].rstrip() if primer_m else ""

    footer_matches = list(FOOTER_RE.finditer(markdown_original))
    if footer_matches:
        linea_inicio = markdown_original.rfind("\n", 0, footer_matches[-1].start()) + 1
        footer = markdown_original[linea_inicio:].strip()
    else:
        footer = ""

    def fmt_h(v: float) -> str:
        return str(int(v)) if v == int(v) else f"{v:.1f}"

    partes: list[str] = []
    if header:
        partes.append(header + "\n\n---\n")

    for bloque in bloques:
        horas_b = float(bloque.get("horas", 0))
        partes.append(
            f"\n## Bloque {bloque['numero']} — {bloque['nombre']} · {fmt_h(horas_b)}h\n\n"
        )
        partes.append("| Subtema | Evidencia | Origen |\n")
        partes.append("|---------|-----------|--------|\n")
        for sub in bloque["subtemas"]:
            if sub.get("manual"):
                evidencia = sub.get("evidencia") or "Manual (profesor)"
                origen = "Manual"
            else:
                evidencia = sub.get("evidencia") or "Sin señal verificable"
                origen = sub.get("origen") or "Detectado"
            partes.append(f"| {sub['nombre']} | {evidencia} | {origen} |\n")
        partes.append("\n")

    partes.append("---\n")
    if footer:
        partes.append(f"\n{footer}\n")

    return "".join(partes)


# ---------------------------------------------------------------------------
# Extracción de horas lectivas y nombre de descarga desde la guía docente.
#
# Lógica pura (solo re + unicodedata, sin Streamlit ni estado). Estas funciones
# son la fuente de verdad importable: tanto el app.py standalone como la
# app-unificada las consumen desde aquí en lugar de duplicarlas.
# ---------------------------------------------------------------------------


def extraer_horas_docencia(texto_guia: str) -> dict[str, int]:
    """
    Extrae horas TE/PA/PL de guías docentes heterogéneas (tablas o texto libre).
    Si no encuentra una categoría, devuelve 0 para esa categoría.
    """
    def norm(texto: str) -> str:
        base = unicodedata.normalize("NFKD", texto or "")
        base = base.encode("ascii", "ignore").decode("ascii").lower()
        return re.sub(r"\s+", " ", base).strip()

    def parsear_numero(token: str) -> int | None:
        try:
            valor = float(token.replace(",", "."))
            if valor < 0:
                return None
            return int(round(valor))
        except ValueError:
            return None

    def categoria_fila_modalidad(linea_norm: str) -> str | None:
        # Reglas mutuamente excluyentes para evitar solapes TE/PA/PL.
        # Incluye raíz "pract" para abreviaturas tipo "Práct. de aula" (sin "practic").
        tiene_practica = ("practic" in linea_norm) or bool(re.search(r"\bpract\b", linea_norm))
        tiene_aula = "aula" in linea_norm
        tiene_seminario = "seminar" in linea_norm
        tiene_taller = "taller" in linea_norm
        tiene_laboratorio = "laborator" in linea_norm or re.search(r"\blab\b", linea_norm) is not None
        tiene_campo = "campo" in linea_norm
        tiene_informatica = "informatica" in linea_norm
        tiene_idiomas = "idioma" in linea_norm
        tiene_teoria = ("expositiv" in linea_norm) or ("teoric" in linea_norm) or ("magistral" in linea_norm)

        # PL primero: categoría más específica y fácil de confundir si no se prioriza.
        if tiene_laboratorio or (tiene_practica and (tiene_campo or tiene_informatica or tiene_idiomas)):
            return "laboratorio"
        # PA: solo con "práctica(s)" (no basta con "seminario"/"taller" sueltos: otras tablas
        # y el texto narrativo también los usan y pueden llevar a 14h en lugar de 7h).
        if tiene_practica and (tiene_aula or tiene_seminario or tiene_taller):
            return "aula"
        # TE: expositiva/teórica/magistral y sin prácticas.
        if tiene_teoria and not tiene_practica:
            return "teoria"
        return None

    def fuerza_fila_pa(linea_norm: str) -> int:
        """Mayor = más seguro que es la fila de modalidad PA (no otra tabla)."""
        if linea_norm is None:
            return 0
        tiene_practica = ("practic" in linea_norm) or bool(re.search(r"\bpract\b", linea_norm))
        if not tiene_practica:
            return 0
        if "laborator" in linea_norm or re.search(r"\blab\b", linea_norm) is not None:
            return 0
        tiene_aula = "aula" in linea_norm
        tiene_seminario = "seminar" in linea_norm
        tiene_taller = "taller" in linea_norm
        tiene_informatica = "informatica" in linea_norm
        tiene_idiomas = "idioma" in linea_norm
        tiene_campo = "campo" in linea_norm
        tiene_laboratorio = "laborator" in linea_norm or re.search(r"\blab\b", linea_norm) is not None
        if tiene_laboratorio or (tiene_practica and (tiene_campo or tiene_informatica or tiene_idiomas)):
            return 0
        if tiene_aula and not (tiene_informatica or tiene_idiomas):
            return 3
        if tiene_seminario or tiene_taller:
            return 2
        return 0

    def es_fila_desglose_temas(linea_norm: str) -> bool:
        """Evita filas de tablas por tema/unidad (suelen repetir columnas TE/PA/PL con otros totales)."""
        if re.match(r"^\d+\s*[\.\)]\s", linea_norm):
            return True
        if re.match(r"^tema\s*\d", linea_norm):
            return True
        if re.match(r"^unidad\s*\d", linea_norm):
            return True
        if re.match(r"^bloque\s*\d", linea_norm):
            return True
        return False

    def es_contexto_horario(linea_norm: str) -> bool:
        # Señales de que el número se refiere a carga horaria y no a conteos (sesiones, grupos...).
        return any(
            patron in linea_norm
            for patron in ("hora", "horas", "h ", " h", "lectiv", "dedicacion", "carga docente", "credit")
        )

    def linea_parece_header_tabla(linea_norm: str) -> bool:
        return (
            ("modalidad" in linea_norm or "modalidades" in linea_norm)
            and ("hora" in linea_norm or "horas" in linea_norm)
        )

    def extraer_hora_fila(linea_original: str) -> int | None:
        nums = [parsear_numero(m.group(0)) for m in re.finditer(r"\d+(?:[.,]\d+)?", linea_original)]
        nums = [n for n in nums if n is not None and 0 <= n <= 500]
        if not nums:
            return None
        candidatos = [n for n in nums if 0 <= n <= 120]
        return candidatos[0] if candidatos else nums[0]

    def es_linea_numerica_de_tabla(linea_original: str, linea_norm: str) -> bool:
        # Acepta líneas "continuación de fila" típicas del OCR de tablas, evita texto narrativo.
        if "sesion" in linea_norm or "sesiones" in linea_norm:
            return False
        numeros = re.findall(r"\d+(?:[.,]\d+)?", linea_original)
        if not numeros:
            return False
        letras = re.findall(r"[A-Za-zÁÉÍÓÚáéíóúÑñ]", linea_original)
        # Heurística: línea casi numérica (poca letra) y al menos 1 número.
        return len(letras) <= 3

    lineas_originales = [l.strip() for l in (texto_guia or "").splitlines() if l.strip()]
    lineas_norm = [norm(l) for l in lineas_originales]

    def extraer_primera_hora_fila_tabla(linea_original: str) -> int | None:
        """
        En tablas MODALIDADES el orden suele ser Horas | % | Totales.
        No usar la posición del texto 'Horas' en la cabecera: cada fila tiene
        prefijo de longitud distinta y el número más cercano puede ser el %.
        """
        coincidencias = list(re.finditer(r"\d+(?:[.,]\d+)?", linea_original))
        for m in coincidencias:
            v = parsear_numero(m.group(0))
            if v is None:
                continue
            if 1900 <= v <= 2100:
                continue
            if 0 <= v <= 120:
                return v
        return None

    horas = {"teoria": 0, "aula": 0, "laboratorio": 0}

    # Estrategia 1 (preferente): detectar tabla(s) MODALIDADES y leer columna Horas.
    # Puede haber varias tablas (p. ej. desglose por temas + resumen); se elige la ventana
    # con más datos y mayor confianza en la fila PA.
    indices_headers = [i for i, ln in enumerate(lineas_norm) if linea_parece_header_tabla(ln)]
    mejor_ventana = None  # (puntuacion, indice_cabecera, dict horas)

    for idx_header in indices_headers:
        local = {"teoria": 0, "aula": 0, "laboratorio": 0}
        fuerza_pa_mejor = -1
        ventana_fin = min(len(lineas_originales), idx_header + 45)

        for i in range(idx_header + 1, ventana_fin):
            ln = lineas_norm[i]
            if es_fila_desglose_temas(ln):
                continue
            if "total" in ln and "modalidad" not in ln:
                continue

            categoria = categoria_fila_modalidad(ln)
            if categoria is None:
                continue

            if categoria == "aula":
                fuerza = fuerza_fila_pa(ln)
                if fuerza < fuerza_pa_mejor:
                    continue

            if categoria != "aula" and local[categoria] != 0:
                continue

            valor = extraer_primera_hora_fila_tabla(lineas_originales[i])

            if valor is None:
                for salto in (1, 2):
                    if i + salto >= ventana_fin:
                        break
                    ln_sig = lineas_norm[i + salto]
                    if es_fila_desglose_temas(ln_sig):
                        continue
                    if categoria_fila_modalidad(ln_sig) is not None or linea_parece_header_tabla(ln_sig):
                        break
                    if not es_linea_numerica_de_tabla(lineas_originales[i + salto], ln_sig):
                        continue
                    valor = extraer_primera_hora_fila_tabla(lineas_originales[i + salto])
                    if valor is not None:
                        break

            if valor is None:
                continue

            if categoria == "aula":
                fuerza = fuerza_fila_pa(ln)
                if fuerza < fuerza_pa_mejor:
                    continue
                local["aula"] = int(valor)
                fuerza_pa_mejor = fuerza
            else:
                local[categoria] = int(valor)

        completitud = sum(1 for v in local.values() if v > 0)
        puntuacion = completitud * 100 + fuerza_pa_mejor
        if mejor_ventana is None:
            mejor_ventana = (puntuacion, idx_header, local)
        elif puntuacion > mejor_ventana[0] or (
            puntuacion == mejor_ventana[0] and idx_header > mejor_ventana[1]
        ):
            # Empate: preferir tabla más abajo en el documento (suele ser el resumen MODALIDADES).
            mejor_ventana = (puntuacion, idx_header, local)

    if mejor_ventana is not None:
        horas = mejor_ventana[2]

    # Estrategia 2 (fallback): búsqueda en texto libre, siempre en la MISMA línea.
    # Se aplica tras tabla y exige señales horarias para evitar capturar "sesiones".
    if 0 in horas.values():
        for linea_original, linea_norm_item in zip(lineas_originales, lineas_norm):
            categoria = categoria_fila_modalidad(linea_norm_item)
            if categoria is None or horas[categoria] != 0:
                continue
            if not es_contexto_horario(linea_norm_item):
                # Evita confundir número de sesiones con horas.
                continue
            valor = extraer_hora_fila(linea_original)
            if valor is not None:
                horas[categoria] = int(valor)

    return {
        "horas_teoria": horas["teoria"],
        "horas_aula": horas["aula"],
        "horas_laboratorio": horas["laboratorio"],
    }


def construir_nombre_descarga(texto_guia: str) -> str:
    """Deriva un nombre de archivo de descarga a partir del nombre de la asignatura."""
    # Se busca el patrón habitual "NOMBRE ... CÓDIGO" para extraer la asignatura.
    coincidencia = re.search(r"NOMBRE\s+(.+?)\s+C[ÓO]DIGO", texto_guia or "", flags=re.IGNORECASE | re.DOTALL)
    if not coincidencia:
        return "Propuesta_asignatura.md"

    nombre_asignatura = coincidencia.group(1).strip()
    # Se eliminan caracteres no válidos para nombres de archivo.
    nombre_limpio = re.sub(r'[\\/:\*\?"<>\|]', "", nombre_asignatura)
    # Se normalizan espacios a guiones bajos.
    nombre_limpio = re.sub(r"\s+", "_", nombre_limpio).strip("_")
    if not nombre_limpio:
        return "Propuesta_asignatura.md"

    return f"Propuesta_{nombre_limpio}.md"


# ---------------------------------------------------------------------------
# Normalización de horas y conteo de bloques sobre el Markdown de organización.
# Lógica pura: transforman/leen el Markdown generado, sin tocar estado de UI.
# ---------------------------------------------------------------------------


def normalizar_horas_output(markdown: str, total_horas: float) -> tuple[str, dict | None]:
    """
    Verifica y normaliza la suma de horas a nivel de BLOQUE en el markdown generado.

    Solo opera sobre los encabezados ## Bloque N — Nombre · Xh. Las horas por subtema
    ya no forman parte del formato de salida.

    Si sum(horas_bloques) != total_horas, redistribuye proporcionalmente (redondeo a
    0.5h). Si la suma ya es correcta, devuelve el markdown sin modificar y None.

    Returns:
        (markdown_corregido, info_ajuste) si se aplicó corrección
        (markdown, None) si la suma era exacta
    """
    if not total_horas or total_horas <= 0:
        return markdown, None

    HDR_RE = re.compile(r"^(## Bloque \d+ — .+? · )([\d.,]+)(h.*)$")

    lineas = markdown.splitlines(keepends=True)
    indices_hdrs: list[tuple[int, float, re.Match]] = []

    for i, linea in enumerate(lineas):
        m = HDR_RE.match(linea.rstrip("\r\n"))
        if not m:
            continue
        try:
            h = float(m.group(2).replace(",", "."))
        except ValueError:
            continue
        indices_hdrs.append((i, h, m))

    if not indices_hdrs:
        return markdown, None

    horas_originales = [h for _, h, _ in indices_hdrs]
    suma_actual = sum(horas_originales)

    if abs(suma_actual - total_horas) < 0.01:
        return markdown, None

    diferencia = suma_actual - total_horas

    def fmt(v: float) -> str:
        return str(int(v)) if v == int(v) else f"{v:.1f}"

    if suma_actual > 0:
        factor = total_horas / suma_actual
        horas_nuevas = [round(h * factor * 2) / 2 for h in horas_originales]
    else:
        n = len(horas_originales)
        base = round((total_horas / n) * 2) / 2 if n else 0
        horas_nuevas = [base] * n

    # Ajuste fino para cuadrar el total exacto.
    ajuste_residuo = total_horas - sum(horas_nuevas)
    if horas_nuevas and abs(ajuste_residuo) >= 0.01:
        horas_nuevas[-1] = round((horas_nuevas[-1] + ajuste_residuo) * 2) / 2

    ajustes: dict[str, dict] = {}
    nuevas_lineas = list(lineas)

    for (idx_linea, h_orig, m), h_new in zip(indices_hdrs, horas_nuevas):
        if abs(h_orig - h_new) >= 0.01:
            nombre_bloque = re.search(r"## Bloque \d+ — (.+?) ·", m.group(0))
            clave = nombre_bloque.group(1).strip() if nombre_bloque else f"Bloque {idx_linea}"
            ajustes[clave] = {"antes": h_orig, "despues": h_new}
        ending = lineas[idx_linea][len(lineas[idx_linea].rstrip("\r\n")):]
        nuevas_lineas[idx_linea] = m.group(1) + fmt(h_new) + m.group(3) + ending

    return "".join(nuevas_lineas), {
        "diferencia": diferencia,
        "suma_antes": suma_actual,
        "ajustes": ajustes,
    }


def contar_bloques_output(markdown_text: str) -> int:
    """Cuenta encabezados ## Bloque N en el markdown generado.

    La regex es estricta: solo cuenta líneas que coinciden exactamente con
    el patrón de bloque temático (## Bloque <número>). Otros encabezados ##
    (subtítulos, secciones internas) no producen falsos positivos.
    """
    return len(re.findall(r"^## Bloque\s+\d+", markdown_text, re.MULTILINE))


# ---------------------------------------------------------------------------
# Parseo enriquecido del Markdown de organización para persistencia (BD).
#
# A diferencia de parsear_bloques_desde_markdown (orientado a la edición manual,
# devuelve solo nombre/horas/manual), esta función captura los metadatos que la
# app-unificada necesita para guardar en la base de datos: orden, evidencia,
# origen y la marca de fallback. La consume la app-unificada al confirmar la
# organización como definitiva.
# ---------------------------------------------------------------------------

# Cabeceras de tabla que no son filas de datos
_CABECERAS_TABLA = {
    "subtema", "topic", "horas", "hours",
    "justificación", "justificacion", "justification",
    "origen", "origin",
    "evidencia", "evidence",
}

# Contrato de formato del Agente Organizador (ver CLAUDE.md del monorepo):
#   ## Bloque N — NOMBRE · Xh
_HDR_BLOQUE_RE = re.compile(
    r"^##\s+Bloque\s+(\d+)\s+[—\-]+\s+(.+?)\s*[·•]\s*([\d.,]+)h",
    re.MULTILINE,
)
# Captura todas las celdas de una fila de tabla Markdown (separadas por '|').
_FILA_TABLA_COMPLETA_RE = re.compile(r"^\|(.+)\|$", re.MULTILINE)


def _parsear_celdas(linea: str) -> list[str]:
    """Extrae celdas de una línea de tabla Markdown como lista de strings limpios."""
    partes = linea.strip().strip("|").split("|")
    return [p.strip() for p in partes]


def parsear_bloques_organizador(markdown: str) -> list[dict]:
    """Extrae la lista de bloques y subtemas del output del Agente Organizador.

    Compatible con los formatos de tabla del Organizador:
      - Actual: | Subtema | Evidencia | Origen |
      - Legado: | Subtema | Horas | Evidencia | Origen |
      - Legado: | Subtema | Horas | Justificación |

    Returns:
        list[dict] con claves: numero (int), nombre (str), horas (float),
        subtemas (list[dict{nombre, horas, orden, evidencia, origen, es_fallback}])
        — horas del subtema se deja en 0.0 (solo aplica a nivel de bloque).
    """
    bloques: list[dict] = []
    matches = list(_HDR_BLOQUE_RE.finditer(markdown))

    for i, m in enumerate(matches):
        numero = int(m.group(1))
        nombre = m.group(2).strip()
        horas_bloque = float(m.group(3).replace(",", "."))

        inicio = m.end()
        fin = matches[i + 1].start() if i + 1 < len(matches) else len(markdown)
        seccion = markdown[inicio:fin]

        subtemas: list[dict] = []
        for fm in _FILA_TABLA_COMPLETA_RE.finditer(seccion):
            celdas = _parsear_celdas(fm.group(0))
            if len(celdas) < 2:
                continue
            col1 = celdas[0]
            col2 = celdas[1] if len(celdas) > 1 else ""
            if col1.lower() in _CABECERAS_TABLA:
                continue
            if re.match(r"^-+$", col1) or re.match(r"^-+$", col2):
                continue
            if not col1:
                continue

            if len(celdas) >= 4 and _es_celda_horas(col2):
                evidencia = celdas[2]
                origen = celdas[3]
            elif len(celdas) >= 4:
                evidencia = celdas[2]
                origen = celdas[3]
            elif len(celdas) == 3 and _es_celda_horas(col2):
                evidencia = ""
                origen = celdas[2]
            elif len(celdas) >= 3:
                evidencia = col2
                origen = celdas[2]
            elif len(celdas) == 2:
                evidencia = col2
                origen = "Detectado"
            else:
                continue

            _ev_norm = evidencia.strip().lower()
            es_fallback = int(
                _ev_norm in {"sin señal verificable", "sin senal verificable", "fallback", ""}
                and origen.strip().lower() in {"fallback", "sin señal", "sin senal", ""}
            )

            subtemas.append({
                "nombre": col1,
                "horas": 0.0,
                "orden": len(subtemas) + 1,
                "evidencia": evidencia,
                "origen": origen,
                "es_fallback": es_fallback,
            })

        bloques.append({
            "numero": numero,
            "nombre": nombre,
            "horas": horas_bloque,
            "subtemas": subtemas,
        })

    return bloques


def parse_organization_md(content: str) -> list[dict]:
    """Adapta `parsear_bloques_organizador` al formato del Agente Contenido.

    Fuente de verdad única: `parsear_bloques_organizador`. Esta función solo
    renombra claves (`subtemas` → `subbloques`) para el standalone de Contenido.
    """
    return [
        {
            "nombre": bloque["nombre"],
            "horas": bloque["horas"],
            "subbloques": [
                {
                    "nombre": subtema["nombre"],
                    "horas": subtema.get("horas", 0.0),
                    "evidencia": subtema.get("evidencia", ""),
                    "origen": subtema.get("origen", ""),
                }
                for subtema in bloque["subtemas"]
            ],
        }
        for bloque in parsear_bloques_organizador(content)
    ]


def detectar_output_truncado(markdown: str, stop_reason: str | None = None) -> dict | None:
    """Detecta si el Markdown de organización quedó incompleto o truncado.

    Señales:
      - stop_reason == 'max_tokens' en la respuesta de la API.
      - Última fila de tabla de subtemas con celdas vacías o incompletas.
      - Último bloque sin tabla de subtemas cerrada correctamente.

    Returns:
        dict con claves 'motivo' (str) y 'detalle' (str), o None si parece íntegro.
    """
    motivos: list[str] = []

    if stop_reason == "max_tokens":
        motivos.append("La respuesta del modelo se cortó por límite de tokens (max_tokens).")

    lineas = [l.strip() for l in (markdown or "").splitlines() if l.strip()]
    ultima_fila_tabla: str | None = None
    cabecera_cols = 0

    for linea in lineas:
        if not linea.startswith("|"):
            continue
        celdas = _parsear_celdas(linea)
        if not celdas:
            continue
        if celdas[0].lower() in _CABECERAS_TABLA or re.match(r"^-+$", celdas[0]):
            cabecera_cols = len(celdas)
            continue
        ultima_fila_tabla = linea

    if ultima_fila_tabla and cabecera_cols >= 3:
        celdas = _parsear_celdas(ultima_fila_tabla)
        nombre_corto = (celdas[0][:80] + "…") if celdas and len(celdas[0]) > 80 else (celdas[0] if celdas else "")
        if len(celdas) < cabecera_cols:
            motivos.append(
                f"Última fila de subtemas incompleta ({len(celdas)} de {cabecera_cols} columnas): "
                f"«{nombre_corto}»"
            )
        elif cabecera_cols == 3 and len(celdas) >= 3:
            if celdas[0] and (not celdas[1].strip() or not celdas[2].strip()):
                motivos.append(f"Fila de subtema sin evidencia u origen: «{nombre_corto}»")
        elif cabecera_cols >= 4 and len(celdas) >= 4:
            if celdas[0] and (not celdas[2].strip() or not celdas[3].strip()):
                motivos.append(f"Fila de subtema sin evidencia u origen: «{nombre_corto}»")

    if not motivos:
        return None

    return {
        "motivo": motivos[0],
        "detalle": " ".join(motivos),
        "stop_reason": stop_reason or "",
    }
